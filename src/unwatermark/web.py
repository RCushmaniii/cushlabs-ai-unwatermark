"""FastAPI web interface — upload, annotate, analyze, remove, compare."""

from __future__ import annotations

import io
import json
import logging
import tempfile
import threading
import uuid
from pathlib import Path

from fastapi import FastAPI, File, Form, Request, UploadFile
from fastapi.exceptions import HTTPException
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, Response, StreamingResponse
from PIL import Image

from unwatermark.cli import _get_handler
from unwatermark.config import load_config
from unwatermark.core.detector import detect_watermark
from unwatermark.models.analysis import WatermarkRegion
from unwatermark.models.annotation import UserAnnotation
from unwatermark.pages import (
    APP_PAGE,
    CONTACT_PAGE,
    HELP_PAGE,
    LANDING_PAGE,
    NOT_FOUND_PAGE,
    PRIVACY_PAGE,
    TERMS_PAGE,
)

logger = logging.getLogger(__name__)

app = FastAPI(title="Unwatermark", description="AI-powered watermark removal")

# Inline SVG favicon — eye-off icon matching the brand
FAVICON_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 32 32">'
    '<rect width="32" height="32" rx="6" fill="#2563eb"/>'
    '<g transform="translate(4,4)" stroke="#fff" stroke-width="1.8" '
    'stroke-linecap="round" fill="none">'
    '<path d="M2 12s3-7 10-7 10 7 10 7-3 7-10 7S2 12 2 12z"/>'
    '<circle cx="12" cy="12" r="3"/>'
    '<line x1="2" y1="2" x2="22" y2="22"/>'
    '</g></svg>'
)

# Download token store — persists to disk so tokens survive server reloads
_DOWNLOAD_DIR = Path(tempfile.gettempdir()) / "unwatermark_downloads"
_DOWNLOAD_DIR.mkdir(exist_ok=True)


# ---------------------------------------------------------------------------
# Pages
# ---------------------------------------------------------------------------

@app.get("/", response_class=HTMLResponse)
async def index():
    return LANDING_PAGE


@app.get("/app", response_class=HTMLResponse)
async def app_page():
    return APP_PAGE


@app.get("/help", response_class=HTMLResponse)
async def help_page():
    return HELP_PAGE


@app.get("/contact", response_class=HTMLResponse)
async def contact_page():
    return CONTACT_PAGE


@app.get("/terms", response_class=HTMLResponse)
async def terms_page():
    return TERMS_PAGE


@app.get("/privacy", response_class=HTMLResponse)
async def privacy_page():
    return PRIVACY_PAGE


@app.get("/favicon.ico")
async def favicon():
    return Response(content=FAVICON_SVG, media_type="image/svg+xml")


@app.get("/healthz")
async def healthz():
    """Health check endpoint — used by external ping services to keep the server warm."""
    return JSONResponse({"status": "ok"})


# ---------------------------------------------------------------------------
# SEO — robots.txt and sitemap.xml
# ---------------------------------------------------------------------------

_SITE_URL = "https://unwatermark.cushlabs.ai"

_ROBOTS_TXT = f"""User-agent: *
Allow: /
Disallow: /analyze
Disallow: /process
Disallow: /download/

Sitemap: {_SITE_URL}/sitemap.xml
"""

_SITEMAP_XML = f"""<?xml version="1.0" encoding="UTF-8"?>
<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">
  <url><loc>{_SITE_URL}/</loc><priority>1.0</priority><changefreq>weekly</changefreq></url>
  <url><loc>{_SITE_URL}/app</loc><priority>0.9</priority><changefreq>weekly</changefreq></url>
  <url><loc>{_SITE_URL}/help</loc><priority>0.7</priority><changefreq>monthly</changefreq></url>
  <url><loc>{_SITE_URL}/contact</loc><priority>0.5</priority><changefreq>monthly</changefreq></url>
  <url><loc>{_SITE_URL}/terms</loc><priority>0.3</priority><changefreq>yearly</changefreq></url>
  <url><loc>{_SITE_URL}/privacy</loc><priority>0.3</priority><changefreq>yearly</changefreq></url>
</urlset>
"""


@app.get("/robots.txt")
async def robots_txt():
    return Response(content=_ROBOTS_TXT.strip(), media_type="text/plain")


@app.get("/sitemap.xml")
async def sitemap_xml():
    return Response(content=_SITEMAP_XML.strip(), media_type="application/xml")


# ---------------------------------------------------------------------------
# API endpoints
# ---------------------------------------------------------------------------

@app.post("/analyze")
async def analyze_file(
    file: UploadFile = File(...),
    description: str = Form(""),
    location: str = Form(""),
    region_x: int = Form(-1),
    region_y: int = Form(-1),
    region_w: int = Form(-1),
    region_h: int = Form(-1),
):
    """Analyze an uploaded file for watermarks."""
    content = await file.read()
    suffix = Path(file.filename).suffix.lower() if file.filename else ""
    is_image = suffix in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}

    # Combine description + location into a single annotation string
    combined_desc = _combine_description(description, location)

    annotation = _build_annotation(
        combined_desc, region_x, region_y, region_w, region_h
    )
    config = load_config()

    if not is_image:
        # Extract first page/slide as an image and run real AI analysis on it
        preview_image = _extract_preview(content, suffix)
        if preview_image is not None:
            try:
                analysis = detect_watermark(preview_image, config, annotation)
                return JSONResponse({
                    "watermark_found": analysis.watermark_found,
                    "region": {
                        "x": analysis.region.x,
                        "y": analysis.region.y,
                        "width": analysis.region.width,
                        "height": analysis.region.height,
                    },
                    "description": analysis.description,
                    "background_type": analysis.background_type.value,
                    "strategy": analysis.strategy.value,
                    "confidence": analysis.confidence,
                    "reasoning": (
                        f"{analysis.reasoning} "
                        "Each page will be analyzed independently during processing."
                    ),
                })
            except Exception as e:
                logger.warning(f"Preview analysis failed, using fallback: {e}")

        # Fallback if preview extraction fails
        return JSONResponse({
            "watermark_found": True,
            "region": {"x": 0, "y": 0, "width": 0, "height": 0},
            "description": combined_desc or "Watermark",
            "background_type": "mixed",
            "strategy": "clone_stamp",
            "confidence": 0.5,
            "reasoning": (
                "Could not preview this document type. "
                "Each page will be analyzed during processing. "
                "Click Remove Watermark to proceed."
            ),
        })

    try:
        image = Image.open(io.BytesIO(content))
    except Exception:
        return JSONResponse(
            {"error": "Could not open file as an image. Check the file format."},
            status_code=400,
        )

    try:
        analysis = detect_watermark(image, config, annotation)
    except Exception as e:
        logger.exception("Analysis failed")
        return JSONResponse(
            {"error": _friendly_error(e)}, status_code=500
        )

    return JSONResponse({
        "watermark_found": analysis.watermark_found,
        "region": {
            "x": analysis.region.x,
            "y": analysis.region.y,
            "width": analysis.region.width,
            "height": analysis.region.height,
        },
        "description": analysis.description,
        "background_type": analysis.background_type.value,
        "strategy": analysis.strategy.value,
        "confidence": analysis.confidence,
        "reasoning": analysis.reasoning,
    })


@app.post("/process")
async def process_file(
    file: UploadFile = File(...),
    description: str = Form(""),
    location: str = Form(""),
    strategy: str = Form(""),
    region_x: int = Form(-1),
    region_y: int = Form(-1),
    region_w: int = Form(-1),
    region_h: int = Form(-1),
):
    """Process with streaming progress — returns newline-delimited JSON events.

    Each line is a JSON object:
      {"type": "progress", "message": "...", "pct": 50}
      {"type": "complete", "download_token": "abc123"}
      {"type": "error", "message": "..."}
    """
    suffix = Path(file.filename).suffix.lower()
    handler = _get_handler(suffix)
    if handler is None:
        return JSONResponse(
            {"error": f"Unsupported file type: {suffix}"},
            status_code=400,
        )

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp_in:
        content = await file.read()
        tmp_in.write(content)
        tmp_in.flush()
        input_path = Path(tmp_in.name)

    combined_desc = _combine_description(description, location)
    annotation = _build_annotation(
        combined_desc, region_x, region_y, region_w, region_h
    )
    config = load_config()
    force_strategy = strategy if strategy else None
    output_path = input_path.with_stem(input_path.stem + "_clean")
    original_filename = file.filename or "file"

    # Progress state shared between the processing thread and the generator
    progress_events: list[dict] = []
    done_event = threading.Event()
    error_holder: list[str] = []

    def on_progress(message: str, pct: int) -> None:
        progress_events.append({"type": "progress", "message": message, "pct": pct})

    def run_handler() -> None:
        try:
            logger.info(
                f"Starting handler: input={input_path}, output={output_path}, "
                f"suffix={suffix}, strategy={force_strategy}"
            )
            handler(
                input_path, output_path, config, annotation, force_strategy,
                on_progress=on_progress,
            )
            logger.info(f"Handler completed. Output exists: {output_path.exists()}")
        except Exception as e:
            logger.exception("Processing failed")
            error_holder.append(_friendly_error(e))
        finally:
            done_event.set()

    # Run the handler in a thread so we can stream progress
    thread = threading.Thread(target=run_handler, daemon=True)
    thread.start()

    def event_stream():
        sent = 0
        while not done_event.is_set():
            done_event.wait(timeout=0.3)
            while sent < len(progress_events):
                yield json.dumps(progress_events[sent]) + "\n"
                sent += 1

        # Flush remaining
        while sent < len(progress_events):
            yield json.dumps(progress_events[sent]) + "\n"
            sent += 1

        if error_holder:
            yield json.dumps({"type": "error", "message": error_holder[0]}) + "\n"
        else:
            token = str(uuid.uuid4())[:8]
            # Persist token mapping to disk so it survives server reloads
            meta = {"path": str(output_path), "name": f"clean_{original_filename}"}
            (_DOWNLOAD_DIR / f"{token}.json").write_text(json.dumps(meta))
            logger.info(f"Download ready: token={token}, path={output_path}")
            yield json.dumps({"type": "complete", "download_token": token}) + "\n"

    return StreamingResponse(event_stream(), media_type="application/x-ndjson")


@app.get("/download/{token}")
async def download_file(token: str):
    """Download a processed file by token."""
    meta_file = _DOWNLOAD_DIR / f"{token}.json"
    if not meta_file.exists():
        logger.error(f"Token {token} not found — meta file missing at {meta_file}")
        return JSONResponse({"error": "Download expired or not found."}, status_code=404)

    meta = json.loads(meta_file.read_text())
    path = Path(meta["path"])
    name = meta.get("name", "clean_file")

    if not path.exists():
        logger.error(f"Output file does not exist: {path}")
        meta_file.unlink(missing_ok=True)
        return JSONResponse({"error": "Download expired or not found."}, status_code=404)

    logger.info(f"Serving download: {path} ({path.stat().st_size} bytes) as {name}")
    # Clean up meta file after serving
    meta_file.unlink(missing_ok=True)
    return FileResponse(
        path=str(path), filename=name, media_type="application/octet-stream"
    )


# ---------------------------------------------------------------------------
# 404 catch-all
# ---------------------------------------------------------------------------

@app.exception_handler(404)
async def not_found_handler(request: Request, exc: HTTPException):
    return HTMLResponse(NOT_FOUND_PAGE, status_code=404)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _combine_description(description: str, location: str) -> str:
    """Merge 'what it looks like' and 'where it is' into a single string."""
    parts = []
    if description.strip():
        parts.append(description.strip())
    if location.strip():
        parts.append(location.strip())
    return ", ".join(parts)


def _build_annotation(
    description: str,
    region_x: int,
    region_y: int,
    region_w: int,
    region_h: int,
) -> UserAnnotation | None:
    """Build a UserAnnotation from form fields, or None if empty."""
    if not description and region_x < 0:
        return None
    region = None
    if region_x >= 0 and region_w > 0:
        region = WatermarkRegion(
            x=region_x, y=region_y, width=region_w, height=region_h
        )
    return UserAnnotation(description=description, region=region)


def _extract_preview(content: bytes, suffix: str) -> Image.Image | None:
    """Extract the first page/slide from a document as a PIL Image for analysis."""
    try:
        if suffix == ".pptx":
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            if not prs.slides:
                return None
            # Find the first image in the first slide
            slide = prs.slides[0]
            for shape in slide.shapes:
                if not shape.shape_type or not hasattr(shape, "image"):
                    continue
                try:
                    image_part = shape.image
                    return Image.open(io.BytesIO(image_part.blob))
                except Exception:
                    continue
            return None

        elif suffix == ".pdf":
            import fitz

            doc = fitz.open(stream=content, filetype="pdf")
            if len(doc) == 0:
                doc.close()
                return None
            page = doc[0]
            # Render at 150 DPI for a quick preview
            zoom = 150 / 72.0
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            doc.close()
            return image

    except Exception as e:
        logger.warning(f"Failed to extract preview from {suffix}: {e}")
    return None


def _friendly_error(exc: Exception) -> str:
    """Convert exceptions into user-friendly error messages."""
    msg = str(exc).lower()
    if "api_key" in msg or "authentication" in msg or "401" in msg:
        return "API key is missing or invalid. Check your .env file."
    if "timeout" in msg or "timed out" in msg:
        return (
            "The AI service took too long to respond. Please try again."
        )
    if "rate" in msg and "limit" in msg:
        return "API rate limit reached. Wait a moment and try again."
    if "connection" in msg:
        return (
            "Could not connect to the AI service. Check your internet "
            "connection and try again."
        )
    return f"Something went wrong: {exc}"
