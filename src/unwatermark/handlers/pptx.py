"""PPTX file handler — processes PowerPoint files with embedded watermarked images.

Targets the common case: slides that are each a single full-size PNG image
(e.g., exported from NotebookLM) with a watermark baked into the image.
"""

from __future__ import annotations

import gc
import io
import logging
from pathlib import Path
from typing import Callable

from PIL import Image
from pptx import Presentation

from unwatermark.config import Config
from unwatermark.core.multipass import clean_image
from unwatermark.models.analysis import WatermarkAnalysis
from unwatermark.models.annotation import UserAnnotation

logger = logging.getLogger(__name__)


def process_pptx(
    input_path: Path,
    output_path: Path,
    config: Config,
    annotation: UserAnnotation | None = None,
    force_strategy: str | None = None,
    on_progress: Callable[[str, int], None] | None = None,
    preview_dir: Path | None = None,
) -> Path:
    """Remove watermarks from all images embedded in a PPTX file.

    Uses multi-pass cleaning on each slide to catch multiple watermarks.
    Tracks a baseline detection from the first successful slide to help
    with slides where detection fails.

    Args:
        input_path: Path to the source PPTX.
        output_path: Path to write the cleaned PPTX.
        config: Runtime configuration.
        annotation: Optional user hints about the watermark.
        force_strategy: Override the AI's strategy recommendation.
        on_progress: Callback(message, percent) for progress updates.

    Returns:
        Path to the output file.
    """
    MAX_SLIDES = 20

    prs = Presentation(str(input_path))

    slide_count = len(prs.slides)
    if slide_count > MAX_SLIDES:
        raise ValueError(
            f"PPTX has {slide_count} slides (max {MAX_SLIDES}). "
            f"Split the file or use the CLI for larger presentations."
        )

    # Baseline from first successful detection — reused when detection fails
    baseline_analysis: WatermarkAnalysis | None = None

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        # Each slide gets a slice of the 5-93% progress range
        slide_start_pct = int(5 + (slide_idx / slide_count) * 88)
        slide_end_pct = int(5 + ((slide_idx + 1) / slide_count) * 88)

        if on_progress:
            on_progress(
                f"Slide {slide_num}/{slide_count}: scanning for watermarks...",
                slide_start_pct,
            )

        for shape in slide.shapes:
            if not shape.shape_type or not hasattr(shape, "image"):
                continue

            try:
                image_wrapper = shape.image
            except Exception:
                continue

            image_bytes = image_wrapper.blob
            content_type = image_wrapper.content_type
            image = Image.open(io.BytesIO(image_bytes))
            logger.info(
                f"Slide {slide_num}: found image {image.width}x{image.height} "
                f"({content_type})"
            )

            # Save "before" preview thumbnail for comparison UI
            if preview_dir is not None:
                _save_preview(image, preview_dir / f"before_{slide_idx}.jpg")

            # Create a sub-progress callback that prefixes messages with slide info
            # and maps clean_image's 10-95% range into this slide's slice
            def _make_slide_progress(s_num: int, s_total: int, start: int, end: int):
                def _slide_progress(msg: str, inner_pct: int) -> None:
                    if on_progress:
                        # Map inner_pct (10-95) into our slide's range (start-end)
                        scaled = start + int((inner_pct - 10) / 85 * (end - start))
                        scaled = max(start, min(end, scaled))
                        on_progress(f"Slide {s_num}/{s_total}: {msg}", scaled)
                return _slide_progress

            slide_progress = _make_slide_progress(slide_num, slide_count, slide_start_pct, slide_end_pct) if on_progress else None

            # Multi-pass: detect and remove all watermarks on this slide
            result = clean_image(
                image, config, annotation, force_strategy,
                baseline=baseline_analysis,
                on_progress=slide_progress,
            )

            if result.removed == 0:
                logger.info(f"Slide {slide_num}: no watermarks removed")
                # Save identical "after" preview for clean slides
                if preview_dir is not None:
                    _save_preview(image, preview_dir / f"after_{slide_idx}.jpg")
                if on_progress:
                    on_progress(f"Slide {slide_num}/{slide_count}: clean", slide_end_pct)
                continue

            logger.info(
                f"Slide {slide_num}: removed {result.removed} watermark(s)"
            )

            if on_progress:
                on_progress(
                    f"Slide {slide_num}/{slide_count}: removed {result.removed} watermark{'s' if result.removed != 1 else ''}",
                    slide_end_pct,
                )

            # Save first successful detection as baseline for future slides
            if baseline_analysis is None and result.first_analysis is not None:
                baseline_analysis = result.first_analysis
                r = baseline_analysis.region
                logger.info(
                    f"Slide {slide_num}: saved baseline "
                    f"({r.x},{r.y},{r.width}x{r.height})"
                )

            cleaned = result.image

            # Save "after" preview thumbnail for comparison UI
            if preview_dir is not None:
                _save_preview(cleaned, preview_dir / f"after_{slide_idx}.jpg")

            buf = io.BytesIO()
            img_format = _content_type_to_format(content_type)
            if img_format == "JPEG":
                cleaned = cleaned.convert("RGB")
            cleaned.save(buf, format=img_format, quality=95)

            # CRITICAL: shape.image returns a wrapper object — setting _blob on it
            # does NOT persist. We must update the actual ImagePart in the package
            # by navigating through the slide's relationship to the embedded image.
            actual_part = _get_image_part(shape, slide)
            if actual_part is not None:
                new_blob = buf.getvalue()
                actual_part._blob = new_blob
                logger.info(
                    f"Slide {slide_num}: replaced blob ({len(new_blob)} bytes)"
                )
            else:
                logger.warning(f"Slide {slide_num}: image part not found")

            # Free per-slide memory — prevents accumulation across 20 slides
            del buf, cleaned, result, image, image_bytes
            gc.collect()

    if on_progress:
        on_progress("Saving presentation...", 95)

    prs.save(str(output_path))

    if on_progress:
        on_progress("Done", 100)

    return output_path


def _get_image_part(shape, slide):
    """Get the actual ImagePart from the PPTX package for a picture shape.

    shape.image returns a read-only wrapper — to persist changes, we need
    the real ImagePart stored in the slide's relationships.
    """
    try:
        pic = shape._element
        ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        blips = pic.findall(f".//{ns_a}blip")
        if not blips:
            return None
        r_embed = blips[0].get(f"{ns_r}embed")
        if not r_embed:
            return None
        return slide.part.related_part(r_embed)
    except Exception as e:
        logger.warning(f"Failed to resolve image part: {e}")
        return None


def _save_preview(image: Image.Image, path: Path, max_width: int = 1200) -> None:
    """Save a JPEG preview thumbnail for the comparison UI."""
    img = image
    if img.width > max_width:
        ratio = max_width / img.width
        img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
    img.convert("RGB").save(path, format="JPEG", quality=85)


def _content_type_to_format(content_type: str) -> str:
    mapping = {
        "image/png": "PNG",
        "image/jpeg": "JPEG",
        "image/gif": "GIF",
        "image/bmp": "BMP",
        "image/tiff": "TIFF",
    }
    return mapping.get(content_type, "PNG")
