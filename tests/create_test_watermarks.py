"""Create a test PPTX with watermarks in varying positions and styles.

Takes the 4-slide test_wm_removal.pptx as baseline, extracts each slide's
image, adds a different watermark to each, and saves a new PPTX.

Watermark layout:
  Slide 1: "Shutterstock" — bottom-right (classic stock photo position)
  Slide 2: "DRAFT" — top-left corner
  Slide 3: "© 2024 Getty Images" — center-bottom
  Slide 4: "SAMPLE PREVIEW" — diagonal across the middle (large, faint)
"""

import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


INPUT = Path.home() / "Downloads" / "test_wm_removal.pptx"
OUTPUT = Path.home() / "Downloads" / "test_varied_watermarks.pptx"

# Watermark definitions: (text, position_func, font_size, color, opacity, rotation)
WATERMARKS = [
    {
        "text": "Shutterstock",
        "position": "bottom-right",
        "font_size": 28,
        "color": (255, 255, 255),
        "opacity": 160,
        "rotation": 0,
    },
    {
        "text": "DRAFT",
        "position": "top-left",
        "font_size": 36,
        "color": (200, 50, 50),
        "opacity": 140,
        "rotation": 0,
    },
    {
        "text": "\u00a9 2024 Getty Images",
        "position": "center-bottom",
        "font_size": 24,
        "color": (255, 255, 255),
        "opacity": 130,
        "rotation": 0,
    },
    {
        "text": "SAMPLE PREVIEW",
        "position": "center",
        "font_size": 72,
        "color": (180, 180, 180),
        "opacity": 80,
        "rotation": -30,
    },
]


def _get_position(pos_name: str, img_w: int, img_h: int, text_w: int, text_h: int):
    """Calculate x, y for a named position."""
    margin = 20
    positions = {
        "bottom-right": (img_w - text_w - margin, img_h - text_h - margin),
        "top-left": (margin, margin),
        "center-bottom": ((img_w - text_w) // 2, img_h - text_h - margin),
        "center": ((img_w - text_w) // 2, (img_h - text_h) // 2),
    }
    return positions[pos_name]


def add_watermark(image: Image.Image, wm: dict) -> Image.Image:
    """Overlay a text watermark on the image."""
    img = image.convert("RGBA")

    # Create transparent overlay for the watermark
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    # Try to use a nice font, fall back to default
    font_size = wm["font_size"]
    try:
        font = ImageFont.truetype("arial.ttf", font_size)
    except OSError:
        try:
            font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", font_size)
        except OSError:
            font = ImageFont.load_default()

    # Measure text
    text = wm["text"]
    bbox = draw.textbbox((0, 0), text, font=font)
    text_w = bbox[2] - bbox[0]
    text_h = bbox[3] - bbox[1]

    # Handle rotation for diagonal watermarks
    if wm["rotation"] != 0:
        # Draw on a temporary image, rotate, then paste
        temp_size = max(text_w, text_h) * 3
        temp = Image.new("RGBA", (temp_size, temp_size), (0, 0, 0, 0))
        temp_draw = ImageDraw.Draw(temp)
        tx = (temp_size - text_w) // 2
        ty = (temp_size - text_h) // 2
        r, g, b = wm["color"]
        temp_draw.text((tx, ty), text, font=font, fill=(r, g, b, wm["opacity"]))
        temp = temp.rotate(wm["rotation"], expand=False, resample=Image.BICUBIC)

        # Center the rotated text on the image
        paste_x = (img.size[0] - temp_size) // 2
        paste_y = (img.size[1] - temp_size) // 2
        overlay.paste(temp, (paste_x, paste_y), temp)
    else:
        x, y = _get_position(wm["position"], img.size[0], img.size[1], text_w, text_h)
        r, g, b = wm["color"]
        draw.text((x, y), text, font=font, fill=(r, g, b, wm["opacity"]))

    # Composite
    result = Image.alpha_composite(img, overlay)
    return result.convert("RGB")


def main():
    prs = Presentation(str(INPUT))
    print(f"Input: {INPUT.name} ({len(prs.slides)} slides)")

    for slide_idx, slide in enumerate(prs.slides):
        wm = WATERMARKS[slide_idx]
        print(f"\nSlide {slide_idx + 1}: '{wm['text']}' at {wm['position']}")

        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            try:
                image_wrapper = shape.image
            except Exception:
                continue

            # Load original image
            image = Image.open(io.BytesIO(image_wrapper.blob))
            print(f"  Original: {image.width}x{image.height}")

            # Add watermark
            watermarked = add_watermark(image, wm)

            # Save watermarked image back into the PPTX
            buf = io.BytesIO()
            watermarked.save(buf, format="PNG", quality=95)

            # Navigate to the actual ImagePart
            pic = shape._element
            ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            blips = pic.findall(f".//{ns_a}blip")
            if blips:
                r_embed = blips[0].get(f"{ns_r}embed")
                if r_embed:
                    part = slide.part.related_part(r_embed)
                    part._blob = buf.getvalue()
                    print(f"  Watermarked: '{wm['text']}' applied")

            # Also save standalone images for visual inspection
            out_dir = Path("test_output")
            out_dir.mkdir(exist_ok=True)
            watermarked.save(str(out_dir / f"slide{slide_idx + 1}_watermarked.png"))
            print(f"  Saved: test_output/slide{slide_idx + 1}_watermarked.png")
            break

    prs.save(str(OUTPUT))
    print(f"\nOutput: {OUTPUT}")
    print("Done! Open the PPTX or check test_output/ for individual slide images.")


if __name__ == "__main__":
    main()
